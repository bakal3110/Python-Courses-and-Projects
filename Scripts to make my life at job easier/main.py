import pyautogui
import keyboard
from time import sleep
from pptx import Presentation
from pptx.util import Inches

# Initialize presentation
prs = Presentation()
current_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Start with an initial slide
slide_counter = 1

# Some coordinates for mouse cursor placement
x_play = 1745
y_play = 335

x_play_adjusted = 1768
y_play_adjusted = 258

x_credentials_name = 859
y_credentials_name = 606
x_credentials_password = 859
y_credentials_password = 656

# Credentials
secret_user = "secret"
secret_user1 = "secret"

secretest_user = "secret"
secretest_user1 = "secret"

def capture_screenshot(filename='screenshot.png', region=None):
    try:
        screenshot = pyautogui.screenshot(region=region) if region else pyautogui.screenshot()
        screenshot.save(filename)
        return filename
    except Exception as e:
        print("Error while taking a screenshot:", e)
        return None

def add_screenshot_to_slide(slide, image_path, width, height, center=False):
    try:
        if center:
            slide_width = prs.slide_width
            left = (slide_width - width) / 2  # Centralizacja
        else:
            left = Inches(0)

        slide.shapes.add_picture(image_path, left, Inches(0), width=width, height=height)
        #print("SS added into slide")
    except Exception as e:
        print("Error while adding a screenshot to slide:", e)

def on_full_screen_hotkey():
    global current_slide
    global slide_counter
    region = (0, 100, 1880, 930)  # Nowy region
    image_path = capture_screenshot(region=region)
    if image_path:
        current_slide = prs.slides.add_slide(prs.slide_layouts[5])
        width = Inches(1880 / 96)
        height = Inches(930 / 96)
        add_screenshot_to_slide(current_slide, image_path, width, height)
        prs.save('presentation.pptx')
        slide_counter += 1
        print(f"\n\tWorkflow SS added to slide {slide_counter}")


def on_partial_screen_hotkey():
    global current_slide
    global slide_counter
    region = (105, 230, 1450, 75)
    
    # open parameters
    pyautogui.moveTo(1692, 339, duration=0)
    pyautogui.click()
    sleep(0.35)
    pyautogui.moveTo(1692, 577, duration=0)
    pyautogui.click()
    sleep(0.35)
    
    image_path = capture_screenshot(region=region)
    if image_path:
        width = Inches(1450 / 96)
        height = Inches(75 / 96)
        add_screenshot_to_slide(current_slide, image_path, width, height, center=False)
        prs.save('presentation.pptx')
        print(f"\t> Infobox SS added on top of slide {slide_counter}")
    
    pyautogui.moveTo(1768, 143, duration=0)
    pyautogui.click()
    
def on_partial_screen_hotkey2():
    global current_slide
    global slide_counter
    region = (105, 230, 1450, 75)
    
    # open parameters
    pyautogui.moveTo(1692, 339, duration=0)
    pyautogui.click()
    sleep(0.35)
    pyautogui.moveTo(1692, 577, duration=0)
    pyautogui.click()
    sleep(0.35)
    
    image_path = capture_screenshot(region=region)
    if image_path:
        width = Inches(1450 / 96)
        height = Inches(75 / 96)
        add_screenshot_to_slide(current_slide, image_path, width, height, center=False)
        prs.save('presentation.pptx')
        print(f"\t> Infobox SS added on top of slide {slide_counter}")
    
    pyautogui.moveTo(1768, 143, duration=0)
    pyautogui.click()
    
def image_to_text():
    region = (1285, 328, 176, 15)
    
    image_path = capture_screenshot(filename="barcodeValue.jpg", region=region)
    #print(image_path) # returns barcodeValue.jpg
    


def pressPlay():
    pyautogui.moveTo(x_play, y_play, duration=0)
    pyautogui.click()
    
def pressPlay_secret():
    pyautogui.moveTo(x_play_adjusted, y_play_adjusted, duration=0)
    pyautogui.click()
    
def secret_user():
    keyboard.press_and_release('backspace')
    pyautogui.moveTo(x_credentials_name, y_credentials_name, duration=0)
    pyautogui.click()
    keyboard.write(secret_user)
    pyautogui.moveTo(x_credentials_password, y_credentials_password, duration=0)
    pyautogui.click()
    sleep(0.3)
    keyboard.write(secret_user1)
    keyboard.press_and_release('enter')

def secretest_user():
    keyboard.press_and_release('backspace')
    pyautogui.moveTo(x_credentials_name, y_credentials_name, duration=0)
    pyautogui.click()
    keyboard.write(secretest_user)
    pyautogui.moveTo(x_credentials_password, y_credentials_password, duration=0)
    pyautogui.click()
    sleep(0.3)
    keyboard.write(secretest_user1)
    keyboard.press_and_release('enter')

def inputNA():
    keyboard.press_and_release('backspace')
    keyboard.write("N/A")
    
def inputPrinter():
    keyboard.press_and_release('backspace')
    keyboard.write("-PHIL171L-SLP")
    
def printPosition():
    print(pyautogui.position())
    
def dispense():
    # select module
    pyautogui.moveTo(860, 415, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 963, duration=0)
    pyautogui.click()
    sleep(0.3)
    # next
    pyautogui.moveTo(1272, 778, duration=0)
    pyautogui.click()
    sleep(2)
    # container type
    pyautogui.moveTo(860,482, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 706, duration=0)
    pyautogui.click()
    sleep(0.3)
    # printer
    pyautogui.moveTo(860, 520, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 569, duration=0)
    pyautogui.click()
    sleep(0.3)
    # next
    pyautogui.moveTo(1286, 778, duration=0)
    pyautogui.click()
    sleep(2)
    # select text box and write 'test'
    pyautogui.moveTo(884, 524, duration=0)
    pyautogui.click()
    keyboard.write("test")
    # press ok
    pyautogui.moveTo(1081, 771, duration=0)
    pyautogui.click()
    sleep(2)
    keyboard.write(secret_user1)
    keyboard.press_and_release('enter')
    sleep(0.3)
    # next
    pyautogui.moveTo(1272, 775, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 461, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 510, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 608, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 635, duration=0)
    pyautogui.click()
    sleep(0.3)
    
    
    # wipe down scrolled down
    pyautogui.moveTo(1359, 645, duration=0)
    pyautogui.click(clicks=4, interval=0.1)
    
    pyautogui.moveTo(860, 496, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 528, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 618, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(860, 640, duration=0)
    pyautogui.click()
    sleep(0.3)
    pyautogui.moveTo(1284, 771, duration=0)
    pyautogui.click()
    
    # select source
    sleep(3)
    pyautogui.moveTo(487, 348, duration=0)
    pyautogui.click()
    sleep(2)
    keyboard.write(secret_user1)
    keyboard.press_and_release('tab')
    keyboard.write(secretest_user)
    keyboard.press_and_release('tab')
    keyboard.write(secretest_user1)
    keyboard.press_and_release('enter')
    pyautogui.moveTo(877, 401, duration=0)
    pyautogui.click()
    
# Workflow Hotkeys assignment
keyboard.add_hotkey('ctrl+shift+1', pressPlay)
keyboard.add_hotkey('ctrl+shift+3', pressPlay_secret)
keyboard.add_hotkey('ctrl+shift+2', secret_user)
keyboard.add_hotkey('ctrl+shift+4', secretest_user)
keyboard.add_hotkey('ctrl+shift+5', inputNA)


# Powerpoint Hotkeys assignment
keyboard.add_hotkey('ctrl+shift+8', on_full_screen_hotkey)
keyboard.add_hotkey('ctrl+shift+9', on_partial_screen_hotkey)
keyboard.add_hotkey('ctrl+shift+7', on_partial_screen_hotkey2)

# Other Hotkeys assignment
keyboard.add_hotkey('ctrl+shift+0', inputPrinter)
keyboard.add_hotkey('ctrl+shift+f8', printPosition)
keyboard.add_hotkey('ctrl+shift+f9', dispense)

# testing hotkeys
keyboard.add_hotkey('ctrl+shift+6', image_to_text)

# Info messages
print("Powerpoint script: ACTIVE")
print("Credentials script: ACTIVE")
print()
print("Powerpoint Hotkeys list:")
print("\t> CTRL+SHIFT+8: SS workflow")
print("\t> CTRL+SHIFT+9: SS ###### data")
print()
print("Workflow Hotkeys list:")
print("\t> CTRL+SHIFT+2: input ###### credentials")
print("\t> CTRL+SHIFT+4: input ###### credentials")
print("\t> CTRL+SHIFT+1: press play button")
print("\t> CTRL+SHIFT+2: press play button (for ######)")
print("\t> CTRL+SHIFT+5: input N/A")
print("\t> CTRL+SHIFT+F9: prepare dispensing")
print()
print("Other:")
print("\t> CTRL+SHIFT+0: print ######")
print("\t> CTRL+SHIFT+F8: Current cursor position (x,y)")
print("\t> CTRL+SHIFT+6: SS Barcode Value in ###### tab")
print()
print("Press ESC to leave script")

keyboard.wait('esc')